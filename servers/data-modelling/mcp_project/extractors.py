"""
Format-specific data extractors for the data-modelling MCP server.

Each public extractor accepts a file path (or raw text) and returns:
    Dict[str, pd.DataFrame]   — table_name → DataFrame

The module also provides:
  normalize_to_frames()      — flatten any dict/list into DataFrames
  detect_and_extract_text()  — smart auto-detect for pasted/inline content
"""

from __future__ import annotations

import io
import json
import re
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import pandas as pd


# ---------------------------------------------------------------------------
# Core normaliser — re-used by JSON, YAML, TOML, XML branches
# ---------------------------------------------------------------------------

def normalize_to_frames(data: Any) -> Dict[str, pd.DataFrame]:
    """
    Recursively flatten a Python object into a dict of DataFrames.

      list of dicts            → {"records": DataFrame}
      dict w/ list-of-dict vals→ one DataFrame per such key
      dict w/ no list vals     → {"data": DataFrame}  (single row)
      plain list of scalars    → {"records": DataFrame(value=...)}
      scalar                   → raises ValueError
    """
    if isinstance(data, list):
        if not data:
            return {}
        if all(isinstance(r, dict) for r in data):
            return {"records": pd.json_normalize(data)}
        return {"records": pd.DataFrame({"value": data})}

    if isinstance(data, dict):
        tables: Dict[str, pd.DataFrame] = {}
        scalars: Dict[str, Any] = {}

        for key, value in data.items():
            safe = re.sub(r"[^A-Za-z0-9_]", "_", str(key))
            if isinstance(value, list) and value and isinstance(value[0], dict):
                tables[safe] = pd.json_normalize(value)
            elif isinstance(value, list) and value:
                tables[safe] = pd.DataFrame({"value": value})
            elif isinstance(value, dict):
                tables[safe] = pd.json_normalize(value)
            else:
                scalars[key] = value

        if not tables:
            # whole dict is one record
            return {"data": pd.json_normalize(data)}

        if scalars:
            tables["_metadata"] = pd.DataFrame([scalars])

        return tables

    raise ValueError(f"Cannot normalise data of type {type(data).__name__}")


# ---------------------------------------------------------------------------
# Word (.docx)
# ---------------------------------------------------------------------------

def extract_docx(file_path: str) -> Dict[str, pd.DataFrame]:
    """
    Extract from a .docx file:
    - Each embedded table  → "table_N"
    - Bullet/numbered lists → "list_N"
    - "Key: Value" paragraphs → "properties"
    - Remaining paragraph text → "content"  (only when nothing else found)
    """
    try:
        from docx import Document
    except ImportError:
        raise ImportError(
            "python-docx is required for .docx files. "
            "Install it with: pip install python-docx"
        )

    doc = Document(file_path)
    result: Dict[str, pd.DataFrame] = {}

    # ---- Embedded tables ----
    for i, tbl in enumerate(doc.tables, 1):
        rows: List[List[str]] = [
            [cell.text.strip() for cell in row.cells]
            for row in tbl.rows
        ]
        if not rows:
            continue
        headers = _dedup_headers(rows[0])
        data_rows = rows[1:]
        df = pd.DataFrame(
            data_rows if data_rows else [[""] * len(headers)],
            columns=headers,
        )
        result[f"table_{i}"] = df

    # ---- Paragraphs ----
    kv: Dict[str, str] = {}
    list_blocks: List[List[str]] = []
    cur_list: List[str] = []
    cur_list_style: str = ""
    content_rows: List[Dict[str, Any]] = []

    def _flush_list() -> None:
        nonlocal cur_list, cur_list_style
        if cur_list:
            n = len(list_blocks) + 1
            key = re.sub(r"[^A-Za-z0-9_]", "_", cur_list_style or "items") or "items"
            list_blocks.append(cur_list[:])
            result[f"{key}_{n}"] = pd.DataFrame({"item": cur_list})
            cur_list = []
            cur_list_style = ""

    for para in doc.paragraphs:
        raw = para.text.strip()
        if not raw:
            _flush_list()
            continue

        style = para.style.name if para.style else ""

        # List paragraphs
        if "list" in style.lower():
            if not cur_list_style:
                cur_list_style = re.sub(r"[^A-Za-z0-9_ ]", "", style).strip()
            cur_list.append(raw)
            continue

        _flush_list()

        # Key: Value  (short key, no newline)
        kv_match = re.match(r"^([^:\n]{1,60}):\s*(.+)$", raw)
        if kv_match:
            kv[kv_match.group(1).strip()] = kv_match.group(2).strip()
            continue

        content_rows.append({"paragraph": len(content_rows) + 1, "style": style, "text": raw})

    _flush_list()

    if kv:
        result["properties"] = pd.DataFrame([kv])

    # Include raw content only when there are few/no structured artefacts
    # or when it's substantial (> 5 paragraphs)
    structural = {k for k in result if not k.startswith("content")}
    if content_rows and (not structural or len(content_rows) > 5):
        result["content"] = pd.DataFrame(content_rows)

    return result


# ---------------------------------------------------------------------------
# PowerPoint (.pptx)
# ---------------------------------------------------------------------------

def extract_pptx(file_path: str) -> Dict[str, pd.DataFrame]:
    """
    Extract from a .pptx file:
    - "slides" table: slide_number, title, body_text, notes
    - Each table shape per slide → "slide_N_table_M"
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError(
            "python-pptx is required for .pptx files. "
            "Install it with: pip install python-pptx"
        )

    prs = Presentation(file_path)
    result: Dict[str, pd.DataFrame] = {}
    slide_rows: List[Dict[str, Any]] = []

    for slide_num, slide in enumerate(prs.slides, 1):
        title_text = ""
        body_parts: List[str] = []
        notes_text = ""
        table_idx = 0

        # Speaker notes
        if slide.has_notes_slide:
            nf = slide.notes_slide.notes_text_frame
            if nf:
                notes_text = nf.text.strip()

        for shape in slide.shapes:
            # ---- Table shapes ----
            if shape.has_table:
                table_idx += 1
                rows: List[List[str]] = [
                    [cell.text.strip() for cell in row.cells]
                    for row in shape.table.rows
                ]
                if rows:
                    headers = _dedup_headers(rows[0])
                    data_rows = rows[1:]
                    result[f"slide_{slide_num}_table_{table_idx}"] = pd.DataFrame(
                        data_rows if data_rows else [[""] * len(headers)],
                        columns=headers,
                    )
                continue

            # ---- Text shapes ----
            if not shape.has_text_frame:
                continue

            is_title = (
                hasattr(shape, "is_placeholder")
                and shape.is_placeholder
                and hasattr(shape, "placeholder_format")
                and shape.placeholder_format is not None
                and shape.placeholder_format.idx == 0
            )
            text = shape.text_frame.text.strip()
            if is_title:
                title_text = text
            elif text:
                body_parts.append(text)

        slide_rows.append(
            {
                "slide_number": slide_num,
                "title": title_text,
                "body_text": "\n".join(body_parts),
                "notes": notes_text,
            }
        )

    if slide_rows:
        result["slides"] = pd.DataFrame(slide_rows)

    return result


# ---------------------------------------------------------------------------
# Markdown (.md / .markdown / inline text)
# ---------------------------------------------------------------------------

def extract_markdown(content: str) -> Dict[str, pd.DataFrame]:
    """
    Extract from Markdown text:
    - YAML front-matter → "frontmatter"
    - Pipe tables       → "table_N"
    - If no tables: heading-structured sections → "sections"
    """
    result: Dict[str, pd.DataFrame] = {}

    # ---- YAML front-matter ----
    fm = re.match(r"^---\n(.*?)\n---", content, re.DOTALL)
    if fm:
        try:
            import yaml
            fm_data = yaml.safe_load(fm.group(1))
            if isinstance(fm_data, dict):
                result["frontmatter"] = pd.DataFrame([fm_data])
        except Exception:
            pass

    # ---- Pipe tables ----
    lines = content.split("\n")
    i = 0
    table_num = 0

    while i < len(lines):
        if lines[i].count("|") >= 2:
            if i + 1 < len(lines) and _is_md_sep(lines[i + 1]):
                headers = _md_row(lines[i])
                rows: List[List[str]] = []
                j = i + 2
                while j < len(lines) and lines[j].count("|") >= 1:
                    row = _md_row(lines[j])
                    row = (row + [""] * len(headers))[: len(headers)]
                    rows.append(row)
                    j += 1
                if rows:
                    table_num += 1
                    result[f"table_{table_num}"] = pd.DataFrame(rows, columns=headers)
                i = j
                continue
        i += 1

    # ---- Sections fallback ----
    if not result:
        sections: List[Dict[str, Any]] = []
        heading = "Introduction"
        level = 0
        body: List[str] = []

        def _flush_section() -> None:
            if body:
                sections.append(
                    {"level": level, "heading": heading, "body": " ".join(body).strip()}
                )
                body.clear()

        for line in lines:
            hm = re.match(r"^(#{1,6})\s+(.+)$", line)
            if hm:
                _flush_section()
                level = len(hm.group(1))
                heading = hm.group(2).strip()
            elif line.strip():
                body.append(line.strip())

        _flush_section()
        if sections:
            result["sections"] = pd.DataFrame(sections)

    return result


def _is_md_sep(line: str) -> bool:
    return bool(re.match(r"^\s*\|?[\s\-:|]+\|[\s\-:|]*$", line))


def _md_row(line: str) -> List[str]:
    line = line.strip().lstrip("|").rstrip("|")
    return [c.strip() for c in line.split("|")]


# ---------------------------------------------------------------------------
# YAML / TOML / XML / HTML
# ---------------------------------------------------------------------------

def extract_yaml(file_path: str) -> Dict[str, pd.DataFrame]:
    """Parse .yaml / .yml and normalise into DataFrames."""
    try:
        import yaml
    except ImportError:
        raise ImportError("pyyaml is required: pip install pyyaml")
    with open(file_path, encoding="utf-8") as fh:
        data = yaml.safe_load(fh)
    return normalize_to_frames(data)


def extract_toml(file_path: str) -> Dict[str, pd.DataFrame]:
    """Parse .toml and normalise into DataFrames."""
    try:
        import tomllib  # Python 3.11+
    except ImportError:
        try:
            import tomli as tomllib  # type: ignore[no-redef]
        except ImportError:
            raise ImportError("tomllib (Python 3.11+) or tomli required for .toml files")
    with open(file_path, "rb") as fh:
        data = tomllib.load(fh)
    return normalize_to_frames(data)


def extract_xml(file_path: str) -> Dict[str, pd.DataFrame]:
    """
    Parse .xml into DataFrames.
    Tries pandas.read_xml first (flat XML); falls back to recursive
    ElementTree → dict → normalize_to_frames for nested XML.
    """
    try:
        df = pd.read_xml(file_path)
        return {"data": df}
    except Exception:
        pass

    tree = ET.parse(file_path)
    root = tree.getroot()
    data = _elem_to_dict(root)
    if not isinstance(data, dict):
        data = {"value": data}
    return normalize_to_frames(data)


def _elem_to_dict(elem: ET.Element) -> Any:
    tag = re.sub(r"\{[^}]+\}", "", elem.tag)  # strip namespace
    children = list(elem)
    d: Dict[str, Any] = {}

    if elem.attrib:
        for k, v in elem.attrib.items():
            k = re.sub(r"\{[^}]+\}", "", k)
            d[f"@{k}"] = v

    if not children:
        text = (elem.text or "").strip()
        if d:
            if text:
                d["#text"] = text
            return d
        return text

    for child in children:
        child_tag = re.sub(r"\{[^}]+\}", "", child.tag)
        val = _elem_to_dict(child)
        if child_tag in d:
            existing = d[child_tag]
            if not isinstance(existing, list):
                d[child_tag] = [existing]
            d[child_tag].append(val)
        else:
            d[child_tag] = val

    if elem.text and elem.text.strip():
        d["#text"] = elem.text.strip()

    return d


def extract_html(file_path: str) -> Dict[str, pd.DataFrame]:
    """Extract all HTML tables from a local file."""
    try:
        tables = pd.read_html(file_path)
        return {f"table_{i + 1}": df for i, df in enumerate(tables)}
    except ImportError:
        raise ImportError("lxml or html5lib required for HTML parsing: pip install lxml")
    except ValueError:
        return {}  # no <table> elements found


# ---------------------------------------------------------------------------
# Smart auto-detector for pasted / inline text
# ---------------------------------------------------------------------------

def detect_and_extract_text(
    text: str,
    hint: str = "auto",
) -> Tuple[Dict[str, pd.DataFrame], str]:
    """
    Detect the format of raw text and extract DataFrames from it.

    Returns (frames, detected_format_name).

    hint options
    ------------
    auto        Try formats in priority order (default)
    json        Parse as JSON
    yaml        Parse as YAML
    csv         Parse as comma-separated
    tsv         Parse as tab-separated
    delimited   Try all common delimiters
    markdown    Parse as Markdown (tables / sections)
    html        Parse as HTML (extract <table> elements)
    key-value   Parse "Key: Value" lines into a properties table
    plain       Store lines as-is in a content table
    """
    text = text.strip()
    if not text:
        return {}, "empty"

    # ---- Explicit hints ----
    if hint == "json":
        return normalize_to_frames(json.loads(text)), "json"

    if hint == "yaml":
        import yaml
        return normalize_to_frames(yaml.safe_load(text)), "yaml"

    if hint in ("csv", "tsv", "delimited"):
        sep = {"csv": ",", "tsv": "\t"}.get(hint)
        df = _try_csv(text, sep=sep)
        if df is not None:
            return {"data": df}, hint
        return {"content": _lines_df(text)}, "plain"

    if hint == "markdown":
        frames = extract_markdown(text)
        return frames, "markdown"

    if hint == "html":
        try:
            tables = pd.read_html(io.StringIO(text))
            return {f"table_{i+1}": df for i, df in enumerate(tables)}, "html"
        except Exception:
            return {}, "html"

    if hint == "key-value":
        frames = _extract_kv(text)
        return frames, "key-value"

    if hint == "plain":
        return {"content": _lines_df(text)}, "plain"

    # ---- Auto-detection (priority order) ----

    # 1. JSON
    if text[:1] in ("{", "["):
        try:
            return normalize_to_frames(json.loads(text)), "json"
        except Exception:
            pass

    # 2. YAML (front-matter or plain key: value document)
    if text.startswith("---") or re.match(r"^[A-Za-z_][\w ]*:\s+\S", text):
        try:
            import yaml
            data = yaml.safe_load(text)
            if isinstance(data, (dict, list)):
                return normalize_to_frames(data), "yaml"
        except Exception:
            pass

    # 3. HTML tables
    if "<table" in text.lower():
        try:
            tables = pd.read_html(io.StringIO(text))
            if tables:
                return {f"table_{i+1}": df for i, df in enumerate(tables)}, "html"
        except Exception:
            pass

    # 4. Markdown tables (need at least 2 pipe-rows)
    if text.count("|") >= 4:
        frames = extract_markdown(text)
        if any(k.startswith("table_") for k in frames):
            return frames, "markdown"

    # 5. Delimited text (CSV/TSV/semicolon/pipe)
    for sep in ("\t", ",", ";"):
        df = _try_csv(text, sep=sep)
        if df is not None:
            fmt = {"t": "tsv", ",": "csv", ";": "csv_semicolon"}.get(sep[0], "csv")
            return {"data": df}, fmt

    # 6. Key: value block
    frames = _extract_kv(text)
    if frames:
        return frames, "key-value"

    # 7. Fallback: line content
    return {"content": _lines_df(text)}, "plain"


# ---------------------------------------------------------------------------
# Text-parsing helpers
# ---------------------------------------------------------------------------

def _try_csv(text: str, sep: Optional[str] = None) -> Optional[pd.DataFrame]:
    kwargs: Dict[str, Any] = {"engine": "python"}
    if sep is None:
        kwargs["sep"] = None  # auto-detect
    else:
        kwargs["sep"] = sep
    try:
        df = pd.read_csv(io.StringIO(text), **kwargs)
        if len(df.columns) >= 2 and len(df) >= 1:
            return df
    except Exception:
        pass
    return None


def _extract_kv(text: str) -> Dict[str, pd.DataFrame]:
    """Parse "Key: Value" or "Key = Value" lines into a properties table."""
    pattern = re.compile(r"^([^:=\n]{1,80})[:\s=]+(.+)$")
    rows: List[Dict[str, str]] = []
    current: Dict[str, str] = {}

    for line in text.split("\n"):
        line = line.strip()
        if not line:
            if current:
                rows.append(current)
                current = {}
            continue
        m = pattern.match(line)
        if m:
            current[m.group(1).strip()] = m.group(2).strip()

    if current:
        rows.append(current)

    return {"properties": pd.DataFrame(rows)} if rows else {}


def _lines_df(text: str) -> pd.DataFrame:
    lines = [ln for ln in text.split("\n") if ln.strip()]
    return pd.DataFrame({"line": range(1, len(lines) + 1), "text": lines})


# ---------------------------------------------------------------------------
# Shared utilities
# ---------------------------------------------------------------------------

def _dedup_headers(headers: List[str]) -> List[str]:
    """Make column headers unique by appending _2, _3, … as needed."""
    seen: Dict[str, int] = {}
    out: List[str] = []
    for h in headers:
        h = h.strip() or "column"
        if h in seen:
            seen[h] += 1
            out.append(f"{h}_{seen[h]}")
        else:
            seen[h] = 1
            out.append(h)
    return out
